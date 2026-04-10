#!/usr/bin/env python3
"""Generate PowerPoint presentation for internship roadmap"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Color scheme
PRIMARY = RGBColor(102, 126, 234)  # Purple-blue
SECONDARY = RGBColor(118, 75, 162)  # Dark purple
TERTIARY = RGBColor(240, 147, 251)  # Light pink
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(245, 247, 250)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = False

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # Add bottom border to title
    title_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    title_shape.line.color.rgb = SECONDARY
    title_shape.line.width = Pt(3)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
add_title_slide(prs, "🏥 CHIVA Ultrasound AI", "Internship Roadmap & System Architecture")

# Slide 2: Overview
add_content_slide(prs, "🎯 Three Core Tasks", [
    "Task 1: LLM + RAG for Shunt Classification",
    "   • Classify shunt types based on escape/re-entry points",
    "   • Ground decisions in CHIVA hemodynamic literature",
    "   • Generate structured clinical reports",
    "",
    "Task 2: LLM + VLM for Active Scanning Guidance",
    "   • Real-time probe positioning recommendations",
    "   • Protocol adherence validation",
    "   • Dynamic context from session findings",
    "",
    "Task 3: Vision-based Vein Identification",
    "   • Real-time vein detection and segmentation",
    "   • Temporal consistency across frames",
    "   • Low-latency processing (<33ms per frame)"
])

# Slide 3: Task 1 Overview
add_content_slide(prs, "📊 Task 1: LLM + RAG for Shunt Classification", [
    "Goal: Automated shunt type identification",
    "",
    "What we're building:",
    "   • Classify shunt types from anatomical findings",
    "   • Ground decisions in CHIVA hemodynamic literature",
    "   • Generate structured clinical reports",
    "   • Minimize hallucination through RAG architecture",
    "   • Benchmark classification against labeled cases",
    "   • Explore fine-tuning strategies if needed"
])

# Slide 4: Task 1 Approach
add_content_slide(prs, "🛠 Task 1: Approach & First Steps", [
    "Phase 1 Approach:",
    "   • LLM Integration: GPT-4/Claude for zero-shot classification",
    "   • RAG Pipeline: CHIVA literature + hemodynamic rules as context",
    "   • Structured Output: JSON schema for consistent results",
    "   • Prompt Engineering: Multi-shot examples + chain-of-thought reasoning",
    "",
    "Phase 2 (if needed):",
    "   • Fine-tuning on labeled CHIVA case studies",
    "   • Domain adaptation using medical literature embeddings",
    "   • Retrieval optimization with semantic chunking",
    "",
    "Tech Stack: OpenAI API, LangChain, FAISS, ChromaDB, Pydantic"
])

# Slide 5: Task 1 Improvements
add_content_slide(prs, "🚀 Task 1: Potential Improvements", [
    "Accuracy Enhancement:",
    "   • Ensemble multiple LLMs • Confidence scoring • Uncertainty quantification",
    "",
    "Latency Optimization:",
    "   • Model distillation • Local model deployment • Caching strategies",
    "",
    "Clinical Integration:",
    "   • Clinician feedback loops • HIPAA-compliant handling • Audit logging",
    "",
    "Robustness:",
    "   • Adversarial testing • Edge case handling • Graceful degradation"
])

# Slide 6: Task 2 Overview
add_content_slide(prs, "🧭 Task 2: Active Scanning Guidance", [
    "Goal: Real-time AI co-pilot for ultrasound operators",
    "",
    "What we're building:",
    "   • Suggest next probe positions based on protocol",
    "   • Validate scanning protocol adherence",
    "   • Update guidance using real-time session findings",
    "   • Maintain low latency (<500ms target)",
    "   • Learn from captured findings to improve guidance",
    "   • Benchmark Vision Language Models for ultrasound"
])

# Slide 7: Task 2 Approach
add_content_slide(prs, "🛠 Task 2: Approach & First Steps", [
    "Phase 1 Approach:",
    "   • LLM Agent: Few-shot prompting with protocol rules",
    "   • Context Management: Session state + captured findings",
    "   • Spatial Reasoning: Rule-based probe guidance logic",
    "   • VLM Survey: Benchmark Vision Transformers for ultrasound",
    "",
    "Phase 2 (if needed):",
    "   • VLM fine-tuning on medical textbook images",
    "   • Multi-modal fusion of image + protocol context",
    "   • Latency-optimized model serving",
    "",
    "Tech Stack: LangChain Agents, CLIP, Redis, FastAPI, vLLM"
])

# Slide 8: Task 2 Improvements
add_content_slide(prs, "🚀 Task 2: Potential Improvements", [
    "Real-Time Performance:",
    "   • Edge inference • Model quantization • Batch optimization",
    "",
    "Multi-Modal Integration:",
    "   • Image feature extraction • Sensor fusion • Temporal consistency",
    "",
    "Operator Experience:",
    "   • Confidence-ranked suggestions • Dismissible recommendations • Preference learning",
    "",
    "Quality Assurance:",
    "   • Protocol coverage metrics • Scanning quality scoring • Skill tracking"
])

# Slide 9: Task 3 Overview
add_content_slide(prs, "🩸 Task 3: Vein Identification & Segmentation", [
    "Goal: Automated vision-based vein detection",
    "",
    "What we're building:",
    "   • Real-time vein detection in ultrasound streams",
    "   • Precise vein segmentation and boundary marking",
    "   • Spatial localization in probe coordinates",
    "   • Low-latency processing (<33ms for 30fps)",
    "   • Temporal consistency across video frames",
    "   • Domain adaptation for robust performance"
])

# Slide 10: Task 3 Approach
add_content_slide(prs, "🛠 Task 3: Approach & First Steps", [
    "Phase 1 Approach:",
    "   • Foundation: Segment Anything Model (SAM) for zero-shot segmentation",
    "   • Detection: Blob detection + morphological operations",
    "   • Enhancement: Edge detection + texture analysis",
    "   • Post-processing: Temporal smoothing for consistency",
    "",
    "Phase 2 (if needed):",
    "   • Fine-tune YOLOv8 on labeled ultrasound datasets",
    "   • Domain adaptation using BUSI dataset",
    "   • Real-time optical flow for temporal tracking",
    "",
    "Tech Stack: SAM (Meta), OpenCV, PyTorch, YOLO v8, scipy"
])

# Slide 11: Task 3 Improvements
add_content_slide(prs, "🚀 Task 3: Potential Improvements", [
    "Detection Accuracy:",
    "   • Fine-tune on ultrasound datasets • Hard example mining • Attention mechanisms",
    "",
    "Speed & Efficiency:",
    "   • Model distillation • ONNX optimization • GPU acceleration",
    "",
    "Temporal Coherence:",
    "   • Optical flow integration • Kalman filtering • Tracking algorithms",
    "",
    "Robustness:",
    "   • Probe angle adaptation • Imaging artifact handling • Multi-scale processing"
])

# Slide 12: Weekly Plan Week 1-2
add_content_slide(prs, "📅 Weekly Plan: Weeks 1-2", [
    "Week 1: Foundation & Setup",
    "   ✓ Project setup & environment configuration",
    "   ✓ Literature review: CHIVA guidelines & hemodynamics",
    "   ✓ RAG pipeline scaffolding with LangChain",
    "   ✓ Benchmark VLM models (CLIP, DINOv2, LLaVA)",
    "   ✓ Review existing vein detection code",
    "",
    "Week 2: Task 1 Foundation (Shunt Classification)",
    "   ✓ Implement zero-shot LLM classification with GPT-4",
    "   ✓ Create CHIVA literature embedding database",
    "   ✓ Design JSON schema for classification output",
    "   ✓ Implement Pydantic models for structured output",
    "   ✓ Build initial prompt templates with chain-of-thought"
])

# Slide 13: Weekly Plan Week 3-4
add_content_slide(prs, "📅 Weekly Plan: Weeks 3-4", [
    "Week 3: Task 1 RAG & Task 2 Start",
    "   ✓ Integrate RAG pipeline with FAISS/ChromaDB",
    "   ✓ Test shunt classification on labeled case studies",
    "   ✓ Implement clinical report generation",
    "   ✓ Start LLM Agent design for guidance",
    "   ✓ Design protocol rule-based system",
    "",
    "Week 4: Task 2 & Task 3 Enhancement",
    "   ✓ Implement real-time guidance LLM agent",
    "   ✓ Integrate session state management (Redis)",
    "   ✓ Optimize vein detection for speed",
    "   ✓ Add temporal smoothing for vein segmentation",
    "   ✓ Performance profiling & latency benchmarking"
])

# Slide 14: Weekly Plan Week 5-6
add_content_slide(prs, "📅 Weekly Plan: Weeks 5-6", [
    "Week 5: Integration & Testing",
    "   ✓ Integrate all three tasks into unified API",
    "   ✓ End-to-end testing with real ultrasound data",
    "   ✓ Benchmark classification accuracy on test set",
    "   ✓ Evaluate guidance quality and adherence",
    "   ✓ Profile latency across all components",
    "",
    "Week 6: Optimization & Documentation",
    "   ✓ Fine-tuning experiments if performance gaps exist",
    "   ✓ Model distillation for inference optimization",
    "   ✓ Comprehensive documentation & API specs",
    "   ✓ Deployment readiness assessment",
    "   ✓ Technical presentation & findings report"
])

# Slide 15: Success Metrics
add_content_slide(prs, "📈 Success Metrics & KPIs", [
    "Task 1: Shunt Classification",
    "   • Classification Accuracy: ≥85% on labeled cases",
    "   • Hallucination Rate: <5% with RAG",
    "   • Latency: <2 seconds per case",
    "   • Report Quality: Clinician review score ≥4/5",
    "",
    "Task 2: Scanning Guidance",
    "   • Protocol Adherence: ≥80% coverage",
    "   • Latency: <500ms per recommendation",
    "   • Operator Acceptance: ≥70% useful suggestions",
    "",
    "Task 3: Vein Detection",
    "   • Detection Sensitivity: ≥90%",
    "   • Segmentation IoU: ≥0.70",
    "   • Real-time: ≥30fps (≤33ms/frame)"
])

# Slide 16: Deployment Strategy
add_content_slide(prs, "🚀 Deployment Strategy", [
    "Deployment Options:",
    "   • Cloud-First: OpenAI API for LLMs, CloudGPU for VLM",
    "   • Hybrid: Local serving for latency-critical tasks",
    "   • Edge: On-device inference using ONNX/TensorRT",
    "",
    "Hardware Requirements:",
    "   • GPU: NVIDIA T4/RTX 3060+ (8GB+ VRAM)",
    "   • CPU: 8-core Intel/AMD processor",
    "   • RAM: 16GB minimum, 32GB recommended",
    "   • Storage: 50GB SSD for models",
    "",
    "Compliance & Security:",
    "   • HIPAA-compliant data handling",
    "   • Audit logging for all decisions",
    "   • Model explainability (SHAP/attention)"
])

# Slide 17: Risk Mitigation
add_content_slide(prs, "⚠️ Risk Mitigation & Contingencies", [
    "Risk: LLM Hallucination",
    "   → Mitigation: RAG pipeline, prompt engineering, confidence thresholding",
    "",
    "Risk: High Latency",
    "   → Mitigation: Optimization, local inference, caching, early exit",
    "",
    "Risk: Low Classification Accuracy",
    "   → Mitigation: Fine-tuning, domain adaptation, ensemble methods",
    "",
    "Risk: Vein Detection Failure",
    "   → Mitigation: Multi-model ensemble, fallback to classical CV",
])

# Slide 18: Technology Stack
add_content_slide(prs, "🔧 Technology Stack Summary", [
    "AI/ML Frameworks:",
    "   • LLMs: OpenAI GPT-4, Anthropic Claude",
    "   • VLMs: CLIP, LLaVA, DINOv2",
    "   • Frameworks: PyTorch, TensorFlow • RAG: LangChain, ChromaDB, FAISS",
    "",
    "Backend & Infrastructure:",
    "   • API: FastAPI • Database: PostgreSQL, Redis",
    "   • Serving: vLLM, TensorRT, ONNX",
    "   • Deployment: Docker, Kubernetes",
    "",
    "Development & Monitoring:",
    "   • Testing: pytest • Monitoring: Prometheus, Grafana",
    "   • Version Control: Git, GitHub"
])

# Slide 19: Deliverables
add_content_slide(prs, "✅ End-of-Internship Deliverables", [
    "Code & Systems:",
    "   ✓ Working prototype with all three tasks integrated",
    "   ✓ Production-ready API with documentation",
    "   ✓ Fine-tuning playbook for future improvements",
    "",
    "Documentation:",
    "   ✓ Comprehensive technical documentation & API specs",
    "   ✓ Deployment guide & infrastructure recommendations",
    "   ✓ Clinical validation results & case studies",
    "",
    "Evaluation:",
    "   ✓ Performance benchmarks & evaluation reports",
    "   ✓ Open-source code repository with examples"
])

# Slide 20: Immediate Next Steps
add_content_slide(prs, "🚀 Immediate Next Steps", [
    "This Week:",
    "   1. Literature review & knowledge base compilation",
    "   2. Environment setup & dependency installation",
    "   3. Initial RAG pipeline scaffolding",
    "   4. VLM benchmark design & experimentation plan",
    "   5. Kick-off meeting with team for feedback",
    "",
    "Next Week:",
    "   • Complete Week 1 deliverables from plan",
    "   • Setup all required data sources and APIs",
    "   • Create baseline models for benchmarking"
])

# Slide 21: Final Slide
add_title_slide(prs, "❓ Questions?", "Let's build the future of ultrasound AI 🚀")

# Save presentation
output_path = '/Users/HariKrishnaD/Downloads/NUS/Internships/Cygnus/cmed_demo/CHIVA_Ultrasound_AI_Internship.pptx'
prs.save(output_path)
print(f"✅ Presentation created: {output_path}")
